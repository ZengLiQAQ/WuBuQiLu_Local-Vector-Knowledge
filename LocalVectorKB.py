
import os
import re
import gc
import time
import hashlib
import markdown
import pytesseract
import numpy as np
import yaml
import torch
from PIL import Image
from loguru import logger
from openpyxl import load_workbook
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer
from bm25_search import BM25Searcher


# ========== 全局配置加载 ==========
def load_config(config_path="config.yaml"):
    """加载配置文件"""
    try:
        with open(config_path, "r", encoding="utf-8") as f:
            config = yaml.safe_load(f)
        return config
    except Exception as e:
        logger.error(f"加载配置文件失败：{e}")
        # 返回默认配置
        return {
            "kb": {"path": "./local_vector_db", "collection_name": "multi_file_kb",
                   "chunk_size": 500, "chunk_overlap": 50, "max_file_size": 500, "incremental_update": True},
            "model": {"embed_model": "BAAI/bge-m3", "device": "auto", "normalize_embeddings": True,
                      "quantization": True, "batch_size": 64},
            "chroma": {"hnsw_space": "cosine", "hnsw_ef_construction": 200, "hnsw_M": 16,
                       "hnsw_ef": 100, "database_impl": "leveldb"},
            "log": {"level": "INFO", "file": "kb_logs.log", "rotation": "100MB", "retention": 7}
        }


# 初始化配置和日志
CONFIG = load_config()
logger.add(
    CONFIG["log"]["file"],
    level=CONFIG["log"]["level"],
    rotation=CONFIG["log"]["rotation"],
    retention=CONFIG["log"]["retention"],
    encoding="utf-8"
)

# ========== MacOS环境适配 ==========
pytesseract.pytesseract.tesseract_cmd = '/usr/local/bin/tesseract'  # MacOS Tesseract路径


class LocalVectorKB:
    """
    本地向量知识库核心类（MacOS优化版）
    支持PDF/MD/Excel/图片/TXT/Word/PPT，MPS加速+量化模型+智能分块
    """
    def __init__(self):
        """初始化（从配置文件加载参数）"""
        # 1. 初始化ChromaDB客户端（LevelDB后端+性能优化）
        self.kb_config = CONFIG["kb"]
        self.chroma_config = CONFIG["chroma"]

        try:
            # ChromaDB 新版本兼容配置（只保留有效参数）
            self.client = chromadb.Client(
                Settings(
                    persist_directory=self.kb_config["path"],
                    anonymized_telemetry=False,
                )
            )
            logger.info(f"ChromaDB初始化成功，存储路径：{self.kb_config['path']}")
        except Exception as e:
            logger.error(f"ChromaDB初始化失败：{e}")
            # 降级到SQLite后端
            self.client = chromadb.Client(
                Settings(
                    persist_directory=self.kb_config["path"],
                    anonymized_telemetry=False,
                )
            )
            logger.warning("降级到SQLite后端运行")

        # 2. 创建/获取集合（余弦相似度）
        self.collection = self.client.get_or_create_collection(
            name=self.kb_config["collection_name"],
            metadata={
                "description": "多文件类型本地向量知识库（MacOS优化）",
                "hnsw:space": self.chroma_config["hnsw_space"],
                "hnsw:construction_ef": self.chroma_config["hnsw_ef_construction"],
                "hnsw:M": self.chroma_config["hnsw_M"]
            }
        )

        # 3. 加载嵌入模型（MPS加速+量化）
        self._init_embed_model()

        # 4. 初始化 BM25 搜索器
        self.bm25_searcher = BM25Searcher()
        self._build_bm25_index()

        # 5. 文件解析器映射
        self.file_parser_map = {
            ".pdf": self._parse_pdf,
            ".md": self._parse_md,
            ".xlsx": self._parse_excel,
            ".jpg": self._parse_image,
            ".png": self._parse_image,
            ".txt": self._parse_txt,
            ".docx": self._parse_docx,
            ".pptx": self._parse_pptx
        }

        logger.info("LocalVectorKB初始化完成")

    def _init_embed_model(self):
        """初始化嵌入模型（MacOS MPS加速+量化）"""
        model_config = CONFIG["model"]

        # 自动检测设备（MPS/CPU）
        if model_config["device"] == "auto":
            device = "mps" if torch.backends.mps.is_available() else "cpu"
        else:
            device = model_config["device"]

        # 禁用量化（MPS不支持）
        use_quantization = model_config.get("quantization", False) and device != "mps"

        # 模型加载参数
        model_kwargs = {}
        if use_quantization:
            model_kwargs["quantization_config"] = {
                "load_in_8bit": True,
            }

        # 加载模型
        try:
            self.embed_model = SentenceTransformer(
                model_config["embed_model"],
                device=device,
                **model_kwargs
            )
            # 编码参数
            self.encode_kwargs = {
                "normalize_embeddings": model_config["normalize_embeddings"],
                "batch_size": model_config["batch_size"],
                "show_progress_bar": False
            }

            # 预加载模型（避免首次调用卡顿）
            self.embed_model.encode(["预加载模型"], **self.encode_kwargs)
            logger.info(f"嵌入模型加载成功：{model_config['embed_model']}（设备：{device}，量化：{use_quantization}）")
        except Exception as e:
            logger.error(f"加载模型失败：{e}")
            # 降级到bge-small-zh-v1.5
            self.embed_model = SentenceTransformer("BAAI/bge-small-zh-v1.5", device=device)
            self.encode_kwargs = {"normalize_embeddings": True, "batch_size": 32, "show_progress_bar": False}
            logger.warning("降级到bge-small-zh-v1.5模型运行")

    def _build_bm25_index(self):
        """构建 BM25 索引"""
        try:
            # 获取所有文档
            all_docs = self.collection.get()
            if all_docs["documents"]:
                documents = [
                    {"id": all_docs["ids"][i], "text": all_docs["documents"][i]}
                    for i in range(len(all_docs["documents"]))
                ]
                self.bm25_searcher.build_index(documents)
                logger.info(f"BM25 索引构建完成，共 {len(documents)} 篇文档")
        except Exception as e:
            logger.error(f"构建 BM25 索引失败：{e}")

    def rebuild_bm25_index(self):
        """重建 BM25 索引"""
        self._build_bm25_index()

    # ========== 私有方法：文件解析 ==========
    def _validate_file(self, file_path):
        """验证文件有效性"""
        # 1. 检查文件是否存在
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在：{file_path}")

        # 2. 检查文件大小
        file_size = os.path.getsize(file_path) / (1024 * 1024)  # MB
        if file_size > self.kb_config["max_file_size"]:
            raise ValueError(f"文件过大（{file_size:.2f}MB），超过限制{self.kb_config['max_file_size']}MB")

        # 3. 检查文件类型
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext not in self.file_parser_map:
            raise ValueError(f"不支持的文件类型：{file_ext}，支持：{list(self.file_parser_map.keys())}")

        return file_ext

    def _parse_pdf(self, file_path):
        """流式解析PDF（减少内存占用）"""
        texts = []
        try:
            with open(file_path, "rb") as f:
                reader = PdfReader(f)
                for page_idx, page in enumerate(reader.pages):
                    # 每解析100页清理一次内存
                    if page_idx % 100 == 0 and page_idx > 0:
                        gc.collect()

                    text = page.extract_text().strip()
                    if text:
                        texts.append(text)
            del reader
            gc.collect()
            logger.info(f"PDF解析完成：{file_path}，共{len(texts)}页")
            return texts
        except Exception as e:
            logger.error(f"PDF解析失败：{file_path}，错误：{e}")
            return []

    def _parse_md(self, file_path):
        """解析Markdown（转为纯文本）"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                md_content = f.read()

            # MD转HTML再转纯文本
            html = markdown.markdown(md_content)
            soup = BeautifulSoup(html, "html.parser")
            text = soup.get_text().strip()

            # 按空行拆分段落
            texts = [t.strip() for t in text.split("\n\n") if t.strip()]
            logger.info(f"MD解析完成：{file_path}，共{len(texts)}段落")
            return texts
        except Exception as e:
            logger.error(f"MD解析失败：{file_path}，错误：{e}")
            return []

    def _parse_excel(self, file_path):
        """解析Excel（按工作表合并）"""
        texts = []
        try:
            wb = load_workbook(file_path, read_only=True)  # 只读模式减少内存
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell else "" for cell in row]
                    row_text = " ".join(row_text).strip()
                    if row_text:
                        sheet_text.append(row_text)
                if sheet_text:
                    texts.append(f"【工作表：{sheet_name}】" + " | ".join(sheet_text))
            wb.close()
            logger.info(f"Excel解析完成：{file_path}，共{len(texts)}工作表")
            return texts
        except Exception as e:
            logger.error(f"Excel解析失败：{file_path}，错误：{e}")
            return []

    def _parse_image(self, file_path):
        """解析图片（OCR识别）"""
        texts = []
        try:
            img = Image.open(file_path)
            # 预处理：转为灰度图提升识别率
            img = img.convert('L')
            text = pytesseract.image_to_string(img, lang="chi_sim+eng").strip()
            if text:
                texts.append(text)
            logger.info(f"图片OCR完成：{file_path}，文本长度：{len(text)}")
            return texts
        except Exception as e:
            logger.error(f"图片解析失败：{file_path}，错误：{e}")
            return []

    def _parse_txt(self, file_path):
        """解析TXT（按行拆分）"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                texts = [t.strip() for t in f.readlines() if t.strip()]
            logger.info(f"TXT解析完成：{file_path}，共{len(texts)}行")
            return texts
        except Exception as e:
            logger.error(f"TXT解析失败：{file_path}，错误：{e}")
            return []

    def _parse_docx(self, file_path):
        """解析Word（段落+表格）"""
        texts = []
        try:
            doc = Document(file_path)
            # 解析段落
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)
            # 解析表格
            for table_idx, table in enumerate(doc.tables):
                table_text = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    row_text = " | ".join(row_text)
                    if row_text:
                        table_text.append(row_text)
                if table_text:
                    texts.append(f"【表格{table_idx+1}】{'; '.join(table_text)}")
            logger.info(f"Word解析完成：{file_path}，共{len(texts)}段内容")
            return texts
        except Exception as e:
            logger.error(f"Word解析失败：{file_path}，错误：{e}")
            return []

    def _parse_pptx(self, file_path):
        """解析PPT（按幻灯片拆分）"""
        texts = []
        try:
            prs = Presentation(file_path)
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                if slide_text:
                    texts.append(f"【幻灯片{slide_idx}】{'; '.join(slide_text)}")
            logger.info(f"PPT解析完成：{file_path}，共{len(texts)}幻灯片")
            return texts
        except Exception as e:
            logger.error(f"PPT解析失败：{file_path}，错误：{e}")
            return []

    def _get_file_hash(self, file_path):
        """生成文件MD5哈希（防重复）"""
        try:
            with open(file_path, "rb") as f:
                return hashlib.md5(f.read()).hexdigest()
        except Exception as e:
            logger.error(f"生成文件哈希失败：{file_path}，错误：{e}")
            return f"error_{time.time()}"

    def _validate_text(self, text):
        """校验文本有效性（过滤乱码/过短文本）"""
        # 1. 过滤空白
        text = text.strip()
        if not text:
            return False

        # 2. 过滤过短文本
        if len(text) < 10:
            logger.debug(f"过滤过短文本：{text[:20]}...")
            return False

        # 3. 过滤乱码（有效字符占比<20%）
        valid_chars = re.findall(r'[\u4e00-\u9fff_a-zA-Z0-9\s，。！？；：""''()（）【】]', text)
        valid_ratio = len(valid_chars) / len(text) if len(text) > 0 else 0
        if valid_ratio < 0.2:
            logger.debug(f"过滤乱码文本（有效占比{valid_ratio:.2f}）：{text[:20]}...")
            return False

        return True

    def _smart_chunk(self, text, file_type):
        """智能分块（按语义+文件类型适配）"""
        # 按文件类型调整分块大小
        type_size_map = {
            ".xlsx": 800, ".pptx": 800, ".docx": 600,
            ".jpg": 300, ".png": 300, ".pdf": 500,
            ".md": 500, ".txt": 400
        }
        chunk_size = type_size_map.get(file_type, self.kb_config["chunk_size"])
        chunk_overlap = self.kb_config["chunk_overlap"]

        # 基于中文标点拆分语义单元
        semantic_separators = "。！？；"
        semantic_units = []
        start = 0
        for i, char in enumerate(text):
            if char in semantic_separators:
                semantic_units.append(text[start:i+1].strip())
                start = i+1
        if start < len(text):
            semantic_units.append(text[start:].strip())

        # 合并为指定大小的块（保留重叠）
        chunks = []
        current_chunk = ""
        for unit in semantic_units:
            if len(current_chunk) + len(unit) <= chunk_size:
                current_chunk += unit
            else:
                if current_chunk:
                    chunks.append(current_chunk)
                    # 保留重叠部分
                    current_chunk = current_chunk[-chunk_overlap:] + unit
                else:
                    # 单个单元超过chunk_size，强制拆分
                    for i in range(0, len(unit), chunk_size - chunk_overlap):
                        chunk = unit[i:i+chunk_size].strip()
                        if chunk:
                            chunks.append(chunk)
                    current_chunk = ""
        if current_chunk and self._validate_text(current_chunk):
            chunks.append(current_chunk)

        # 过滤无效文本
        chunks = [c for c in chunks if self._validate_text(c)]
        return chunks

    # ========== 公有方法：核心功能 ==========
    def add_file(self, file_path):
        """
        添加文件到知识库（增量更新+智能分块）
        :param file_path: 文件路径
        :return: 入库文本块数
        """
        try:
            # 1. 验证文件
            file_ext = self._validate_file(file_path)

            # 2. 增量更新检查
            file_mtime = os.path.getmtime(file_path)
            if self.kb_config["incremental_update"]:
                # 查询现有数据
                existing_data = self.collection.get(where={"file_path": file_path})
                if existing_data["metadatas"]:
                    existing_mtime = max([meta.get("file_mtime", 0) for meta in existing_data["metadatas"]])
                    if file_mtime <= existing_mtime:
                        logger.info(f"文件未修改，跳过入库：{file_path}")
                        return 0
                    # 删除旧数据
                    self.collection.delete(where={"file_path": file_path})
                    logger.info(f"文件已修改，删除旧数据：{file_path}")

            # 3. 解析文件
            raw_texts = self.file_parser_map[file_ext](file_path)
            if not raw_texts:
                logger.warning(f"文件无有效文本：{file_path}")
                return 0

            # 4. 智能分块
            chunks = []
            for text in raw_texts:
                chunks.extend(self._smart_chunk(text, file_ext))
            if not chunks:
                logger.warning(f"分块后无有效文本：{file_path}")
                return 0

            # 5. 生成ID和元数据
            file_hash = self._get_file_hash(file_path)
            ids = [f"{file_hash}_{i}" for i in range(len(chunks))]
            metadatas = [
                {
                    "file_path": file_path,
                    "file_type": file_ext,
                    "chunk_index": i,
                    "file_mtime": file_mtime,
                    "file_hash": file_hash,
                    "version": int(time.time())
                } for i in range(len(chunks))
            ]

            # 6. 批量向量化（MPS加速）
            embeddings = self.embed_model.encode(chunks, **self.encode_kwargs).tolist()

            # 7. 入库
            self.collection.add(
                documents=chunks,
                embeddings=embeddings,
                ids=ids,
                metadatas=metadatas
            )
            # ChromaDB 新版本自动持久化

            logger.info(f"文件入库成功：{file_path}，分块数：{len(chunks)}")
            # 重建 BM25 索引
            if len(chunks) > 0:
                self.rebuild_bm25_index()
            return len(chunks)
        except Exception as e:
            logger.error(f"文件入库失败：{file_path}，错误：{e}")
            return 0

    def add_files(self, file_paths):
        """批量添加文件"""
        total_chunks = 0
        for file_path in file_paths:
            if os.path.exists(file_path):
                total_chunks += self.add_file(file_path)
            else:
                logger.warning(f"文件不存在：{file_path}")
        # 重建 BM25 索引
        if total_chunks > 0:
            self.rebuild_bm25_index()
        logger.info(f"批量入库完成，总块数：{total_chunks}")
        return total_chunks

    def search(self, query, top_k=5, file_type_filter=None, hybrid=True, use_bm25=True):
        """
        混合检索（向量 + BM25）
        :param query: 查询文本
        :param top_k: 返回条数
        :param file_type_filter: 文件类型过滤
        :param hybrid: 是否开启混合检索
        :param use_bm25: 是否使用 BM25
        :return: 格式化结果
        """
        try:
            # 1. 构建过滤条件
            where = None
            if file_type_filter:
                where = {"file_type": file_type_filter}

            # 2. 向量检索
            query_embedding = self.embed_model.encode([query], **self.encode_kwargs).tolist()
            vector_results = self.collection.query(
                query_embeddings=query_embedding,
                n_results=top_k * 2 if hybrid else top_k,
                where=where
            )

            if not vector_results["ids"][0]:
                return []

            # 3. BM25 检索（如果启用）
            bm25_scores = {}
            if use_bm25 and hybrid:
                bm25_results = self.bm25_searcher.search(query, top_k * 2)
                bm25_scores = {doc_id: score for doc_id, score in bm25_results}

            # 4. 构建结果并计算综合得分
            query_words = [w for w in query.strip().split() if len(w) > 1]
            formatted_results = []

            # 向量权重和 BM25 权重
            vector_weight = 0.6
            bm25_weight = 0.4

            # 获取向量检索结果
            vector_dict = {}
            for i in range(len(vector_results["ids"][0])):
                doc_id = vector_results["ids"][0][i]
                vector_dict[doc_id] = {
                    "text": vector_results["documents"][0][i],
                    "distance": vector_results["distances"][0][i],
                    "metadata": vector_results["metadatas"][0][i]
                }

            # 合并向量和 BM25 结果
            all_doc_ids = set(vector_dict.keys()) | set(bm25_scores.keys())

            for doc_id in all_doc_ids:
                if doc_id not in vector_dict:
                    continue

                doc_data = vector_dict[doc_id]
                text = doc_data["text"]
                distance = doc_data["distance"]
                metadata = doc_data["metadata"]

                # 向量得分
                vector_score = 1 / (1 + distance)

                # BM25 得分（归一化）
                bm25_score = 0.0
                if doc_id in bm25_scores:
                    max_bm25 = max(bm25_scores.values()) if bm25_scores else 1
                    bm25_score = bm25_scores[doc_id] / max_bm25 if max_bm25 > 0 else 0

                # 关键词匹配度
                match_count = sum([1 for w in query_words if w in text]) if query_words else 0
                match_score = match_count / len(query_words) if query_words else 0

                # 综合得分
                if hybrid and use_bm25:
                    total_score = vector_score * vector_weight + bm25_score * bm25_weight
                elif hybrid:
                    total_score = vector_score * 0.7 + match_score * 0.3
                else:
                    total_score = vector_score

                # 关键词高亮
                highlighted_text = text
                for word in query_words:
                    if len(word) > 1:
                        highlighted_text = re.sub(f"({word})", r"<mark>\1</mark>", highlighted_text, flags=re.IGNORECASE)

                formatted_results.append({
                    "id": doc_id,
                    "text": text,
                    "highlighted_text": highlighted_text,
                    "similarity_distance": round(distance, 4),
                    "vector_score": round(vector_score, 4),
                    "bm25_score": round(bm25_score, 4),
                    "match_score": round(match_score, 4),
                    "total_score": round(total_score, 4),
                    "metadata": metadata
                })

            # 按综合得分排序
            formatted_results = sorted(formatted_results, key=lambda x: x["total_score"], reverse=True)[:top_k]

            logger.info(f"检索完成：{query}，返回{len(formatted_results)}条结果")
            return formatted_results
        except Exception as e:
            logger.error(f"检索失败：{query}，错误：{e}")
            return []

    def delete_by_file(self, file_path):
        """删除指定文件的所有数据"""
        try:
            self.collection.delete(where={"file_path": file_path})
            # ChromaDB 新版本自动持久化
            logger.info(f"删除文件数据成功：{file_path}")
            return True
        except Exception as e:
            logger.error(f"删除文件数据失败：{file_path}，错误：{e}")
            return False

    def batch_delete(self, file_paths):
        """批量删除文件数据"""
        success_count = 0
        for file_path in file_paths:
            if self.delete_by_file(file_path):
                success_count += 1
        logger.info(f"批量删除完成：成功{success_count}/{len(file_paths)}")
        return success_count

    def get_stats(self):
        """获取知识库统计信息"""
        try:
            # 获取文件类型分布
            all_metadatas = self.collection.get()["metadatas"]
            file_type_dist = {}
            file_list = {}

            if all_metadatas:
                for meta in all_metadatas:
                    # 统计文件类型
                    file_type = meta.get("file_type", "unknown")
                    file_type_dist[file_type] = file_type_dist.get(file_type, 0) + 1

                    # 构建文件列表
                    file_path = meta.get("file_path")
                    if file_path not in file_list:
                        file_list[file_path] = {
                            "path": file_path,
                            "type": file_type,
                            "chunk_count": 1,
                            "mtime": meta.get("file_mtime"),
                            "version": meta.get("version")
                        }
                    else:
                        file_list[file_path]["chunk_count"] += 1

            return {
                "total_chunks": self.collection.count(),
                "collection_name": self.collection.name,
                "storage_path": self.kb_config["path"],
                "supported_file_types": list(self.file_parser_map.keys()),
                "file_type_dist": file_type_dist,
                "file_list": list(file_list.values())
            }
        except Exception as e:
            logger.error(f"获取统计信息失败：{e}")
            return {"total_chunks": 0, "error": str(e)}

    def clear_all(self):
        """清空知识库（谨慎使用）"""
        try:
            self.collection.delete()
            # ChromaDB 新版本自动持久化
            logger.warning("知识库已清空")
            return True
        except Exception as e:
            logger.error(f"清空知识库失败：{e}")
            return False


# ========== 测试代码 ==========
if __name__ == "__main__":
    # 初始化知识库
    kb = LocalVectorKB()

    # 测试统计信息
    stats = kb.get_stats()
    logger.info(f"知识库统计：{stats}")

    # 示例：添加文件
    # kb.add_file("./test.pdf")

    # 示例：检索
    # results = kb.search("测试查询", top_k=3)
    # logger.info(f"检索结果：{results}")
